"""
SSE plc Interview Preparation — PowerPoint Generator
Produces a polished, infographic-rich .pptx file using python-pptx + matplotlib
"""

import io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib.patches as patches

# ── Brand colours ──────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x2D, 0x5B)
GREEN  = RGBColor(0x00, 0xA6, 0x51)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xF4, 0xF6, 0xF9)
DGREY  = RGBColor(0x6B, 0x72, 0x80)
LGREENB= RGBColor(0xE6, 0xF7, 0xEE)

NAVY_HEX  = "#1B2D5B"
GREEN_HEX = "#00A651"
WHITE_HEX = "#FFFFFF"
GREY_HEX  = "#F4F6F9"
LGREY_HEX = "#E8EDF4"

# ── Slide dimensions (widescreen 16:9) ─────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


# ── Low-level drawing helpers ───────────────────────────────────────────────
def rect(slide, l, t, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    return shape


def txt(slide, text, l, t, w, h,
        size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size     = Pt(size)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb = color
    return txb


def add_mpl_figure(slide, fig, l, t, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='none', transparent=True)
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(l), Inches(t), Inches(w), Inches(h))
    plt.close(fig)


def nav_bar(slide, active_idx):
    """Thin green bottom bar + page label"""
    rect(slide, 0, 7.2, 13.33, 0.3, NAVY)
    labels = ["Overview","Divisions","Projects","Investment","Impact","Values","EE Masters","Interview"]
    for i, label in enumerate(labels):
        x = 0.3 + i * 1.6
        c = GREEN if i == active_idx else WHITE
        txt(slide, label, x, 7.22, 1.5, 0.26,
            size=7, bold=(i==active_idx), color=c, align=PP_ALIGN.CENTER)


def section_header(slide, label, title, subtitle=""):
    rect(slide, 0, 0, 13.33, 1.3, NAVY)
    rect(slide, 0, 1.3, 13.33, 0.06, GREEN)
    txt(slide, label.upper(), 0.4, 0.1, 6, 0.35,
        size=8, bold=True, color=GREEN)
    txt(slide, title, 0.4, 0.38, 10, 0.7,
        size=26, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 1.0, 10, 0.35,
            size=11, color=RGBColor(0xCC,0xD6,0xE8))


# ══════════════════════════════════════════════════════════════════════════════
#  CHART HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def make_donut_investment():
    """Donut chart: £33bn allocation"""
    fig, ax = plt.subplots(figsize=(4.5, 4.5))
    fig.patch.set_alpha(0)
    ax.set_facecolor('none')

    sizes  = [22, 5, 4, 2]
    labels = ['SSEN\nTransmission\n£22bn', 'SSEN\nDistribution\n£5bn',
              'SSE\nRenewables\n£4bn', 'SSE\nThermal\n£2bn']
    colors = ['#1B2D5B', '#00A651', '#243a78', '#33c76b']
    explode = (0.04, 0.04, 0.04, 0.04)

    wedges, texts = ax.pie(sizes, colors=colors, explode=explode,
                           startangle=90, wedgeprops=dict(width=0.55))

    ax.text(0, 0, '£33bn\nTotal', ha='center', va='center',
            fontsize=14, fontweight='bold', color='#1B2D5B')

    ax.legend(wedges, labels, loc='lower center', bbox_to_anchor=(0.5, -0.22),
              ncol=2, fontsize=8, frameon=False)
    ax.set_title('5-Year Investment Allocation', fontsize=11,
                 fontweight='bold', color='#1B2D5B', pad=10)
    plt.tight_layout()
    return fig


def make_capacity_bar():
    """Bar chart: renewable capacity growth"""
    fig, ax = plt.subplots(figsize=(5.5, 3.8))
    fig.patch.set_alpha(0)
    ax.set_facecolor('#F4F6F9')

    years  = ['2024\n(Now)', '2027\nTarget', '2032\nTarget']
    values = [4.4, 9, 16]
    colors = ['#6B7280', '#00A651', '#1B2D5B']

    bars = ax.bar(years, values, color=colors, width=0.5, edgecolor='white', linewidth=1.5)
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width()/2., bar.get_height() + 0.3,
                f'{val} GW', ha='center', va='bottom',
                fontsize=12, fontweight='bold', color='#1B2D5B')

    ax.set_ylim(0, 20)
    ax.set_ylabel('GW of Renewable Capacity', fontsize=9, color='#6B7280')
    ax.set_title('SSE Renewables — Capacity Growth Roadmap', fontsize=11,
                 fontweight='bold', color='#1B2D5B', pad=10)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.tick_params(colors='#6B7280')
    ax.yaxis.set_tick_params(labelsize=8)
    plt.tight_layout()
    return fig


def make_impact_bar():
    """Horizontal bar: UK economic impact"""
    fig, ax = plt.subplots(figsize=(5.8, 3.2))
    fig.patch.set_alpha(0)
    ax.set_facecolor('#F4F6F9')

    metrics = ['Jobs\nSupported', 'Scottish Economy\n(£bn)', 'UK GVA\n(£bn)', 'Daily Investment\n(£M)']
    values  = [83, 3.4, 9.7, 7]
    colors  = ['#1B2D5B', '#00A651', '#1B2D5B', '#00A651']

    bars = ax.barh(metrics, values, color=colors, edgecolor='white', height=0.5)
    for bar, val in zip(bars, values):
        label = f'{val}k' if val == 83 else f'£{val}bn' if val > 5 else f'£{val}M/day' if val == 7 else f'£{val}bn'
        ax.text(bar.get_width() + 0.5, bar.get_y() + bar.get_height()/2,
                label, va='center', fontsize=10, fontweight='bold', color='#1B2D5B')

    ax.set_xlim(0, 100)
    ax.set_title('SSE\'s UK Economic Contribution', fontsize=11,
                 fontweight='bold', color='#1B2D5B', pad=10)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_visible(False)
    ax.xaxis.set_visible(False)
    ax.tick_params(axis='y', colors='#1B2D5B', labelsize=9)
    plt.tight_layout()
    return fig


def make_timeline_chart():
    """Scatter/line: SSE investment timeline"""
    fig, ax = plt.subplots(figsize=(7, 2.8))
    fig.patch.set_alpha(0)
    ax.set_facecolor('none')

    years  = [2022, 2024, 2026, 2028, 2030, 2032]
    events = [
        (2022, 'Viking\nConstruction\nBegins'),
        (2024, 'Viking\nOperational\n443MW'),
        (2025, '£33bn Plan\nAnnounced'),
        (2027, '9GW\nRenewables\nTarget'),
        (2029, 'Peterhead\nCCS\nOperational'),
        (2032, '16GW\nRenewables\nTarget'),
    ]
    ex = [e[0] for e in events]
    ey = [1]*len(events)

    ax.plot([2022, 2032], [1, 1], color='#00A651', linewidth=3, zorder=1)
    ax.scatter(ex, ey, s=120, color='#1B2D5B', zorder=3, edgecolors='white', linewidth=2)

    for i, (year, label) in enumerate(events):
        ypos = 1.35 if i % 2 == 0 else 0.58
        ax.text(year, ypos, label, ha='center', va='center' if i%2==0 else 'center',
                fontsize=7, color='#1B2D5B', fontweight='bold',
                bbox=dict(boxstyle='round,pad=0.3', facecolor='#E6F7EE',
                          edgecolor='#00A651', linewidth=0.8))
        ax.plot([year, year], [1.15 if i%2==0 else 0.85,
                               1.28 if i%2==0 else 0.72],
                color='#00A651', linewidth=1, linestyle='--')

    ax.set_xlim(2021, 2033)
    ax.set_ylim(0.3, 1.7)
    ax.axis('off')
    ax.set_title('SSE Strategic Timeline', fontsize=10, fontweight='bold',
                 color='#1B2D5B', pad=6)
    plt.tight_layout()
    return fig


def make_values_radar():
    """Radar/spider chart for SSE values"""
    fig, ax = plt.subplots(figsize=(4, 4), subplot_kw=dict(polar=True))
    fig.patch.set_alpha(0)
    ax.set_facecolor('none')

    labels = ['Safety', 'Service', 'Efficiency', 'Sustainability', 'Excellence', 'Teamwork']
    N = len(labels)
    values = [10, 8, 8, 9, 8, 9]
    values += values[:1]

    angles = [n / float(N) * 2 * np.pi for n in range(N)]
    angles += angles[:1]

    ax.set_theta_offset(np.pi / 2)
    ax.set_theta_direction(-1)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(labels, size=9, fontweight='bold', color='#1B2D5B')
    ax.set_ylim(0, 10)
    ax.set_yticks([])
    ax.plot(angles, values, color='#00A651', linewidth=2)
    ax.fill(angles, values, color='#00A651', alpha=0.25)
    ax.spines['polar'].set_color('#1B2D5B')
    ax.grid(color='#1B2D5B', alpha=0.2)
    ax.set_title('SSE SET Values', fontsize=10, fontweight='bold',
                 color='#1B2D5B', pad=18)
    plt.tight_layout()
    return fig


def make_process_funnel():
    """Funnel/steps chart for recruitment process"""
    fig, ax = plt.subplots(figsize=(5, 4))
    fig.patch.set_alpha(0)
    ax.set_facecolor('none')

    steps = [
        ('1. Online Application', '#1B2D5B', 200),
        ('2. SHL Online Tests\n(SJT + Ability)', '#243a78', 160),
        ('3. Video Interview\n(4Q | 3min think | 4min answer)', '#2d4999', 120),
        ('4. Assessment Centre\n(Group + Presentation + Interview)', '#00A651', 80),
    ]

    for i, (label, color, width) in enumerate(steps):
        y = 3.2 - i * 0.85
        bar = ax.barh(y, width, height=0.6, color=color, left=(200-width)/2)
        ax.text(100, y, label, ha='center', va='center',
                fontsize=8.5, fontweight='bold', color='white')

    ax.set_xlim(0, 200)
    ax.set_ylim(-0.2, 3.9)
    ax.axis('off')
    ax.set_title('SSE Recruitment Funnel', fontsize=10, fontweight='bold',
                 color='#1B2D5B', pad=8)
    plt.tight_layout()
    return fig


def make_ee_skills_chart():
    """Horizontal bars for EE skills relevance"""
    fig, ax = plt.subplots(figsize=(6, 4.2))
    fig.patch.set_alpha(0)
    ax.set_facecolor('#F4F6F9')

    skills = [
        'Power Systems Analysis',
        'HV/MV Electrical Design',
        'Renewable Energy Systems',
        'Power Electronics',
        'Protection & Control',
        'EMC Engineering',
        'Smart Grid & Storage',
        'Safety-Critical Eng.',
    ]
    relevance = [98, 95, 90, 88, 92, 80, 82, 96]
    colors = ['#1B2D5B' if r > 90 else '#00A651' for r in relevance]

    y_pos = np.arange(len(skills))
    bars = ax.barh(y_pos, relevance, color=colors, edgecolor='white', height=0.6)
    for bar, val in zip(bars, relevance):
        ax.text(bar.get_width() + 0.5, bar.get_y() + bar.get_height()/2,
                f'{val}%', va='center', fontsize=8, fontweight='bold', color='#1B2D5B')

    ax.set_yticks(y_pos)
    ax.set_yticklabels(skills, fontsize=8.5)
    ax.set_xlim(0, 110)
    ax.set_xlabel('Relevance at SSE (%)', fontsize=8, color='#6B7280')
    ax.set_title('EE Masters Skills → SSE Role Relevance', fontsize=10,
                 fontweight='bold', color='#1B2D5B', pad=8)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.tick_params(axis='y', colors='#1B2D5B')
    ax.tick_params(axis='x', colors='#6B7280', labelsize=8)

    legend_elements = [mpatches.Patch(color='#1B2D5B', label='Critical (>90%)'),
                       mpatches.Patch(color='#00A651', label='High (80-90%)')]
    ax.legend(handles=legend_elements, loc='lower right', fontsize=7, frameon=True)
    plt.tight_layout()
    return fig


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    sl = blank_slide(prs)
    # Full navy background
    rect(sl, 0, 0, 13.33, 7.5, NAVY)
    # Green accent stripe
    rect(sl, 0, 5.4, 13.33, 0.12, GREEN)
    # Green left bar
    rect(sl, 0, 0, 0.18, 7.5, GREEN)

    # SSE wordmark
    txt(sl, 'SSE', 0.5, 0.5, 4, 1.2, size=72, bold=True, color=WHITE)
    txt(sl, 'Scottish and Southern Energy plc', 0.5, 1.55, 7, 0.4,
        size=12, color=RGBColor(0xAA, 0xBB, 0xDD))

    # Divider
    rect(sl, 0.5, 2.05, 6, 0.04, GREEN)

    # Title
    txt(sl, 'Interview Preparation Guide', 0.5, 2.2, 9, 0.85,
        size=32, bold=True, color=WHITE)
    txt(sl, 'Electrical Engineering Masters Edition', 0.5, 2.95, 9, 0.5,
        size=16, color=GREEN, bold=True)

    # Stat pills (manual boxes)
    pills = [
        ('FTSE 100', 0.5),
        ('HQ: Perth, Scotland', 2.3),
        ('£20.93bn Market Cap', 4.9),
        ('£7M Invested Daily', 7.8),
    ]
    for label, x in pills:
        rect(sl, x, 3.65, len(label)*0.11 + 0.3, 0.38, RGBColor(0x2B,0x45,0x80))
        txt(sl, label, x+0.12, 3.68, len(label)*0.11+0.1, 0.32,
            size=9, bold=True, color=WHITE)

    # Bottom note
    txt(sl, 'Includes real candidate insights from TheStudentRoom · AptitudePrep · SSE Careers',
        0.5, 6.85, 10, 0.4, size=9, color=RGBColor(0x88,0x99,0xBB), italic=True)

    # Right side: mini stat boxes
    stat_data = [('£33bn', '5-Year Investment Plan'), ('83,000+','Jobs Supported'),
                 ('16 GW','Renewables by 2032'), ('25%+','UK Offshore Wind 2030')]
    for i, (val, lbl) in enumerate(stat_data):
        y = 1.2 + i * 1.35
        rect(sl, 10.0, y, 3.0, 1.15, RGBColor(0x22,0x3A,0x70))
        txt(sl, val, 10.1, y+0.08, 2.8, 0.55, size=24, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        txt(sl, lbl, 10.1, y+0.6, 2.8, 0.45, size=9, color=WHITE, align=PP_ALIGN.CENTER)
    return sl


def slide_company_overview(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LGREY)
    section_header(sl, 'Company Overview', 'What is SSE plc?',
                   'A FTSE 100 energy company at the heart of the UK\'s net zero transition')
    nav_bar(sl, 0)

    # Stat boxes row
    stats = [('💷 £20.93bn', 'Market Capitalisation'), ('👷 83,000+', 'Jobs Supported'),
             ('⚡ ~£7M/day', 'Daily Investment'), ('📍 Perth', 'HQ, Scotland')]
    for i, (val, lbl) in enumerate(stats):
        x = 0.4 + i * 3.2
        rect(sl, x, 1.65, 2.9, 1.25, NAVY)
        txt(sl, val, x+0.1, 1.72, 2.7, 0.55, size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        txt(sl, lbl, x+0.1, 2.2, 2.7, 0.45, size=9, color=WHITE, align=PP_ALIGN.CENTER)

    # History text box
    rect(sl, 0.4, 3.1, 5.8, 3.85, WHITE)
    txt(sl, 'History & Background', 0.55, 3.18, 5.5, 0.38, size=13, bold=True, color=NAVY)
    history = (
        "• Formed in 1998 from merger of Scottish Hydro-Electric & Southern Electric\n"
        "• Listed on London Stock Exchange — constituent of FTSE 100\n"
        "• Strategic pivot since 2018: sold household supply to focus entirely on "
        "low-carbon electricity infrastructure\n"
        "• Operates across UK and Ireland with growing international renewables presence\n"
        "• Uniquely positioned as a pure-play energy transition business"
    )
    txt(sl, history, 0.55, 3.62, 5.55, 3.1, size=10, color=NAVY)

    # Mission blockquote
    rect(sl, 6.5, 3.1, 6.5, 1.5, LGREENB)
    rect(sl, 6.5, 3.1, 0.12, 1.5, GREEN)
    txt(sl, '"Providing the energy needed today while\nbuilding a better world of energy for tomorrow."',
        6.75, 3.2, 6.1, 0.95, size=12, italic=True, bold=False, color=NAVY)
    txt(sl, '— SSE plc, Annual Report', 6.75, 4.15, 6.1, 0.35, size=9, color=DGREY, italic=True)

    # Why SSE is different
    rect(sl, 6.5, 4.75, 6.5, 2.2, WHITE)
    txt(sl, 'Why SSE is Different from Other Energy Companies',
        6.65, 4.82, 6.2, 0.4, size=11, bold=True, color=NAVY)
    diff = ("• NOT a consumer-facing energy supplier (sold that business)\n"
            "• 100% focused on infrastructure: grids, wind, flexible generation\n"
            "• One of the UK's largest single private capital investors\n"
            "• Regulated revenue streams provide long-term stability\n"
            "• Engineers work on projects at the frontier of the energy transition")
    txt(sl, diff, 6.65, 5.3, 6.2, 2.1, size=9.5, color=NAVY)
    return sl


def slide_business_divisions(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LGREY)
    section_header(sl, 'Business Divisions', 'How SSE is Structured',
                   'Five principal units — know which one you are interviewing for')
    nav_bar(sl, 1)

    divisions = [
        ('SSE Renewables', 'Onshore/offshore wind, hydro, solar, battery storage',
         '9GW by 2027 → 16GW by 2032\n25%+ of UK 2030 offshore wind target\nDogger Bank: world\'s largest OWF', GREEN),
        ('SSEN Transmission', 'HV electricity transmission network, north Scotland',
         '£22bn investment planned\n11 Pathway to 2030 projects\n20%+ of UK networks investment', NAVY),
        ('SSEN Distribution', 'DNO: south England + north Scotland',
         '3.9M homes & businesses served\n£5bn planned investment\nSmart grid & EV network upgrades', RGBColor(0x1F,0x5C,0x35)),
        ('SSE Thermal', 'Flexible generation: gas, CCS, hydrogen, EfW',
         'Peterhead CCS: Scotland\'s first\n910MW + 1.5Mt CO₂/yr captured\nHydrogen transition ready', RGBColor(0x2D,0x3A,0x70)),
        ('SSE Energy Solutions', 'B2B low-carbon local energy infrastructure',
         'EV charging, heat networks, solar PV\nPrivate wire networks for industry\nEnergy-as-a-service', RGBColor(0x00,0x7A,0x3D)),
    ]

    cols = [(0.25, 2.5), (2.9, 2.5), (5.55, 2.5), (8.2, 2.5), (10.5, 2.75)]
    for i, (title, desc, details, color) in enumerate(divisions):
        x, w = cols[i]
        rect(sl, x, 1.55, w, 5.6, WHITE)
        rect(sl, x, 1.55, w, 0.85, color)
        txt(sl, title, x+0.1, 1.6, w-0.2, 0.4, size=10, bold=True, color=WHITE)
        txt(sl, desc, x+0.1, 2.0, w-0.2, 0.35, size=7.5, color=WHITE, italic=True)
        txt(sl, details, x+0.12, 2.5, w-0.25, 3.5, size=8.5, color=NAVY)
    return sl


def slide_key_projects(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LGREY)
    section_header(sl, 'Key Projects', 'Flagship Programmes You Must Know',
                   'Name these projects in your interview — shows genuine research depth')
    nav_bar(sl, 2)

    projects = [
        ('1', 'Dogger Bank\nWind Farm',
         'WORLD\'S LARGEST OFFSHORE WIND FARM',
         'Phases A+B+C = 3.6GW + Phase D = 1.5GW\nPartners: SSE, Equinor, Vårgrønn\nYorkshire coast, North Sea\nHaliade-X 13MW+ turbines\nFoundations complete: Dec 2025',
         GREEN),
        ('2', 'Viking\nWind Farm',
         'EUROPE\'S LARGEST ONSHORE WIND FARM',
         '443MW onshore — Shetland Islands\nOperational: August 2024\nRequired dedicated HVDC interconnector\n(first ever subsea link to Shetland)\nMajor Scottish energy milestone',
         NAVY),
        ('3', 'Peterhead CCS\nPower Station',
         'SCOTLAND\'S FIRST CCS PLANT',
         '910MW gas + carbon capture\nUp to 1.5M tonnes CO₂/yr by 2030\nSSE Thermal + Equinor partnership\nPart of Acorn Carbon Cluster\nHydrogen economy enabler',
         RGBColor(0x1F,0x5C,0x35)),
        ('4', 'Pathway\nto 2030',
         'SSEN TRANSMISSION GRID PROGRAMME',
         '11 major grid projects in north Scotland\nCore of £22bn transmission investment\nNew OHL, substations, HVDC cables\nEnables renewable power to flow south\nDelivering RIIO-T2 & future RIIO-T3',
         RGBColor(0x2D,0x3A,0x70)),
    ]

    for i, (num, title, badge, details, color) in enumerate(projects):
        x = 0.25 + i * 3.27
        rect(sl, x, 1.55, 3.05, 5.55, WHITE)
        rect(sl, x, 1.55, 3.05, 1.1, color)
        txt(sl, num, x+0.12, 1.58, 0.4, 0.45, size=20, bold=True, color=WHITE)
        txt(sl, title, x+0.5, 1.58, 2.5, 0.7, size=11, bold=True, color=WHITE)
        rect(sl, x+0.12, 2.72, 2.82, 0.3, RGBColor(0xE6,0xF7,0xEE))
        txt(sl, badge, x+0.14, 2.74, 2.8, 0.26, size=6.5, bold=True, color=NAVY)
        txt(sl, details, x+0.12, 3.1, 2.85, 3.8, size=8.5, color=NAVY)
    return sl


def slide_investment(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LGREY)
    section_header(sl, 'Investment & Strategy', 'SSE\'s £33bn Five-Year Investment Plan',
                   'Announced November 2025 — one of the UK\'s largest private infrastructure programmes')
    nav_bar(sl, 3)

    # Donut chart
    fig = make_donut_investment()
    add_mpl_figure(sl, fig, 0.3, 1.6, 5.2, 4.8)

    # Breakdown text
    breakdown = [
        ('£22bn', 'SSEN Transmission', '67% of total — grid upgrades in north Scotland', GREEN),
        ('£5bn',  'SSEN Distribution', '15% of total — smart grid & EV network', NAVY),
        ('£4bn',  'SSE Renewables',    '12% of total — new wind, solar, storage', RGBColor(0x1F,0x5C,0x35)),
        ('£2bn',  'SSE Thermal + Other','6% of total — CCS, hydrogen, flexibility', DGREY),
    ]
    for i, (amount, div, desc, color) in enumerate(breakdown):
        y = 1.65 + i * 1.35
        rect(sl, 5.8, y, 0.1, 1.1, color)
        rect(sl, 5.95, y, 7.1, 1.1, WHITE)
        txt(sl, amount, 6.05, y+0.05, 1.4, 0.55, size=24, bold=True, color=color)
        txt(sl, div, 7.3, y+0.06, 5.5, 0.38, size=11, bold=True, color=NAVY)
        txt(sl, desc, 7.3, y+0.44, 5.5, 0.45, size=9, color=DGREY, italic=True)

    # Key stat
    rect(sl, 5.8, 7.0, 7.25, 0.0, NAVY)
    txt(sl, '80% of the entire plan goes directly into the electricity grid — this is fundamentally a grid engineering story.',
        5.8, 6.85, 7.2, 0.4, size=9, bold=True, color=NAVY, italic=True)
    return sl


def slide_capacity_timeline(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LGREY)
    section_header(sl, 'Investment & Growth', 'Renewable Capacity Roadmap + Strategic Timeline',
                   'SSE is targeting a fivefold increase in renewables output by 2032')
    nav_bar(sl, 3)

    fig1 = make_capacity_bar()
    add_mpl_figure(sl, fig1, 0.3, 1.55, 6.2, 4.2)

    fig2 = make_timeline_chart()
    add_mpl_figure(sl, fig2, 6.7, 1.55, 6.4, 3.2)

    # Key callouts
    callouts = [
        ('Viking Wind Farm\nOperational', '2024', GREEN),
        ('£33bn Plan\nAnnounced', '2025', NAVY),
        ('9GW Target', '2027', GREEN),
        ('Peterhead CCS\nOnline', '~2029', NAVY),
        ('16GW Target', '2032', GREEN),
    ]
    for i, (label, year, color) in enumerate(callouts):
        x = 6.9 + i * 1.28
        rect(sl, x, 4.95, 1.15, 1.05, color)
        txt(sl, year, x+0.05, 4.98, 1.05, 0.32, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, label, x+0.05, 5.32, 1.05, 0.65, size=7.5, color=WHITE, align=PP_ALIGN.CENTER)

    return sl


def slide_uk_impact(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LGREY)
    section_header(sl, 'UK Impact', 'How SSE Has Transformed Britain\'s Energy Landscape',
                   'From infrastructure investment to jobs and economic growth')
    nav_bar(sl, 4)

    fig = make_impact_bar()
    add_mpl_figure(sl, fig, 0.3, 1.55, 6.5, 4.3)

    # Impact cards
    impacts = [
        ('Energy Security', '#1B2D5B',
         'Reduces UK exposure to volatile global energy markets by building homegrown generation & infrastructure'),
        ('Job Creation', '#00A651',
         '83,000+ jobs in engineering, construction, operations and supply chains across Scotland and England'),
        ('Net Zero Leadership', '#1B2D5B',
         'Delivering 25%+ of UK\'s 2030 offshore wind target and 20%+ of all electricity network investment'),
        ('Economic Output', '#00A651',
         '£9.7bn GVA added to UK economy; £3.4bn specifically to Scottish economy annually'),
    ]
    for i, (title, color, desc) in enumerate(impacts):
        row, col = divmod(i, 2)
        x = 7.1 + col * 3.05
        y = 1.58 + row * 2.8
        rect(sl, x, y, 2.9, 2.55, WHITE)
        rect(sl, x, y, 2.9, 0.42, RGBColor(*[int(color[j:j+2], 16) for j in (1,3,5)]))
        txt(sl, title, x+0.1, y+0.06, 2.7, 0.32, size=10, bold=True, color=WHITE)
        txt(sl, desc, x+0.1, y+0.52, 2.7, 1.85, size=8.5, color=NAVY)

    return sl


def slide_values(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LGREY)
    section_header(sl, 'Company Values', 'SSE SET — What SSE Stands For',
                   'These values underpin every competency question. Know them cold.')
    nav_bar(sl, 5)

    fig = make_values_radar()
    add_mpl_figure(sl, fig, 0.3, 1.5, 4.8, 4.8)

    values = [
        ('S', 'Safety', '#1B2D5B', 'TOP PRIORITY',
         'Non-negotiable in HV environments. Mention this FIRST if asked about SSE\'s values. '
         'Demonstrate a safety-first mindset in every competency answer.'),
        ('S', 'Service', '#00A651', 'COMMUNITY FOCUS',
         'Serving communities, customers and stakeholders. SSE operates licensed infrastructure '
         'with obligations to the public.'),
        ('E', 'Efficiency', '#1B2D5B', 'OPERATIONAL EXCELLENCE',
         'Responsible stewardship of capital — critical when spending £7M every day on infrastructure.'),
        ('S', 'Sustainability', '#00A651', 'NET ZERO LEADERSHIP',
         'Every investment decision is judged against long-term environmental impact and 2050 net zero targets.'),
        ('E', 'Excellence', '#1B2D5B', 'ENGINEERING STANDARDS',
         'Setting and maintaining the highest professional standards — directly relevant to your EE Masters.'),
        ('T', 'Teamwork', '#00A651', 'CROSS-DISCIPLINE COLLABORATION',
         'Multidisciplinary teams: engineers, project managers, environmental, commercial, policy.'),
    ]

    col_positions = [(5.3, 1.55), (9.1, 1.55), (5.3, 3.35), (9.1, 3.35), (5.3, 5.15), (9.1, 5.15)]
    for i, (letter, name, color, badge, desc) in enumerate(values):
        x, y = col_positions[i]
        c = RGBColor(*[int(color[j:j+2], 16) for j in (1,3,5)])
        rect(sl, x, y, 3.6, 1.75, WHITE)
        rect(sl, x, y, 0.55, 1.75, c)
        txt(sl, letter, x+0.05, y+0.5, 0.45, 0.7, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, name, x+0.65, y+0.08, 2.85, 0.38, size=12, bold=True, color=NAVY)
        txt(sl, badge, x+0.65, y+0.45, 2.85, 0.28, size=7, bold=True, color=c)
        txt(sl, desc, x+0.65, y+0.72, 2.85, 1.0, size=7.5, color=NAVY)

    rect(sl, 0.3, 6.87, 12.7, 0.0, NAVY)
    txt(sl, '"Doing the Right Thing" — SSE\'s overarching leadership ethos underpinning all six values',
        0.3, 6.82, 12.7, 0.35, size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER, italic=True)
    return sl


def slide_ee_masters(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LGREY)
    section_header(sl, 'Your EE Masters', 'How Your Degree Maps to SSE\'s Needs',
                   'Link every skill to a named project or SSE business unit in your answers')
    nav_bar(sl, 6)

    fig = make_ee_skills_chart()
    add_mpl_figure(sl, fig, 0.3, 1.55, 6.5, 5.1)

    # Table on right
    rows = [
        ('Power Systems Analysis', 'Load flow, fault studies — SSEN Transmission daily'),
        ('HV/MV Electrical Design', '132kV–400kV substations — Pathway to 2030'),
        ('Renewable Energy Systems', 'Turbine integration — SSE Renewables'),
        ('Power Electronics', 'HVDC converters, inverters, battery storage'),
        ('Protection & Control', 'Relay co-ord, SCADA — live HV network'),
        ('EMC Engineering', 'Cable design, earthing in HV environments'),
        ('Smart Grid & Storage', 'Battery management, demand-side response'),
        ('Safety-Critical Eng.', 'IEC/BS standards, RAMS, HV site work'),
    ]
    rect(sl, 7.0, 1.55, 6.1, 5.3, WHITE)
    txt(sl, 'EE Skill → SSE Role Mapping', 7.1, 1.6, 5.9, 0.38, size=11, bold=True, color=NAVY)
    rect(sl, 7.0, 1.95, 6.1, 0.35, NAVY)
    txt(sl, 'SKILL', 7.1, 2.0, 2.8, 0.28, size=8, bold=True, color=WHITE)
    txt(sl, 'SSE RELEVANCE', 9.85, 2.0, 3.2, 0.28, size=8, bold=True, color=WHITE)

    for i, (skill, rel) in enumerate(rows):
        y = 2.35 + i * 0.59
        bg = LGREY if i % 2 == 0 else WHITE
        rect(sl, 7.0, y, 6.1, 0.58, bg)
        txt(sl, skill, 7.08, y+0.08, 2.75, 0.4, size=8.5, bold=True, color=NAVY)
        txt(sl, rel, 9.87, y+0.05, 3.1, 0.48, size=7.8, color=NAVY)

    rect(sl, 7.0, 7.0, 6.1, 0.0, GREEN)
    txt(sl, '💡 Name a project: "My power electronics module covered HVDC — directly relevant to Viking\'s Shetland interconnector."',
        7.0, 6.88, 6.1, 0.38, size=7.5, bold=True, color=NAVY, italic=True)
    return sl


def slide_interview_process(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LGREY)
    section_header(sl, 'Interview Prep', 'The Full SSE Recruitment Process',
                   'Sourced from real candidates on TheStudentRoom & AptitudePrep')
    nav_bar(sl, 7)

    fig = make_process_funnel()
    add_mpl_figure(sl, fig, 0.3, 1.5, 4.5, 5.2)

    steps = [
        ('1', 'Online Application', '#1B2D5B',
         'CV + motivation questions\nTailor to your specific division (Transmission, Renewables, Thermal, Distribution)\n~200 graduates hired per year — competitive but achievable'),
        ('2', 'SHL Online Tests', '#00A651',
         '• SJT: 16 questions (values alignment)\n• Ability Test: 24 numerical/analytical/deductive questions\n• Described by candidates as "absolutely brutal — timed out at Q8"\n• Practice at shl.com — same format as NatWest, KPMG etc.\n• No verbal reasoning for engineering roles'),
        ('3', 'Video Interview', '#1B2D5B',
         '• 4 questions | 3 min think | 4 min answer per question\n• Real questions asked: "Why this division?", "3 energy challenges"\n• Competency: Safety + Initiative\n• Can be pre-recorded or live\n• Professional background, test audio beforehand'),
        ('4', 'Assessment Centre', '#00A651',
         '• Group discussion + individual presentation + role play\n• Presentation brief is "intentionally vague" — structure it yourself\n• Transmission final stage was "more technical than expected"\n• Virtual or in-person — ask recruiter for format in advance\n• ⚠ No visa sponsorship for graduate engineering roles'),
    ]

    for i, (num, title, color, detail) in enumerate(steps):
        c = RGBColor(*[int(color[j:j+2], 16) for j in (1,3,5)])
        y = 1.58 + i * 1.42
        rect(sl, 5.1, y, 0.55, 1.25, c)
        txt(sl, num, 5.1, y+0.3, 0.55, 0.55, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, 5.7, y, 7.4, 1.25, WHITE)
        txt(sl, title, 5.82, y+0.06, 7.1, 0.36, size=12, bold=True, color=NAVY)
        txt(sl, detail, 5.82, y+0.44, 7.1, 0.78, size=8, color=NAVY)
    return sl


def slide_star_questions(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LGREY)
    section_header(sl, 'Interview Prep', 'STAR Method + Key Questions',
                   'SSE officially recommends STAR — prepare your answers before the 3-minute think time')
    nav_bar(sl, 7)

    # STAR box
    rect(sl, 0.3, 1.55, 5.0, 3.0, NAVY)
    txt(sl, 'The STAR Technique', 0.45, 1.62, 4.7, 0.38, size=12, bold=True, color=GREEN)
    star = [('S','Situation','Set the scene — context, when, where'),
            ('T','Task','Your specific responsibility or challenge'),
            ('A','Action','"I" not "we" — your specific decisions'),
            ('R','Result','Outcome — quantify where possible')]
    for i, (letter, word, desc) in enumerate(star):
        y = 2.1 + i * 0.57
        rect(sl, 0.45, y, 0.5, 0.45, GREEN)
        txt(sl, letter, 0.45, y+0.03, 0.5, 0.38, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, word, 1.05, y+0.03, 1.2, 0.38, size=10, bold=True, color=WHITE)
        txt(sl, desc, 2.3, y+0.05, 2.9, 0.35, size=8.5, color=RGBColor(0xCC,0xDD,0xEE))

    # Example
    rect(sl, 0.3, 4.65, 5.0, 2.05, RGBColor(0x22,0x3A,0x70))
    txt(sl, 'EE Masters STAR Example:', 0.45, 4.72, 4.7, 0.32, size=9, bold=True, color=GREEN)
    example = ('"During my Masters dissertation [S], I modelled transient stability of a '
               'grid-connected wind farm [T]. I built a MATLAB/Simulink simulation and '
               'proposed a modified converter control strategy [A], reducing fault settling '
               'time by 23% — submitted to IEEE PES [R]." → Link to SSEN Transmission.')
    txt(sl, example, 0.45, 5.05, 4.7, 1.58, size=8, color=WHITE, italic=True)

    # Common questions
    questions = [
        ('Why SSE / this division?',
         'Reference: net zero mission, specific project (Dogger Bank / Peterhead / Pathway to 2030), EE skills match, £33bn scale'),
        ('3 energy challenges facing SSE?',
         'Thermal: grid intermittency, CCS economics, hydrogen transition\nTransmission: N-S bottlenecks, supply chain lead times, planning'),
        ('Time you demonstrated SAFETY?',
         'HV lab risk assessment, stopping unsafe work, pre-task briefing — link to SSE\'s #1 value explicitly'),
        ('Complex engineering problem solved?',
         'Final-year project, simulation, design trade-off — explain thinking process, not just the answer'),
        ('Example of teamwork?',
         'Group design project, cross-discipline collaboration, adapting communication style for non-specialists'),
        ('Energy sector in 10 years?',
         'Offshore wind dominance, grid transformation (SSE\'s 80% bet), hydrogen, digitalisation, storage revolution'),
        ('Questions to ask SSE?',
         '"What would you expect from me in the first 90 days?" — directly recommended by SSE\'s own recruiter blog'),
    ]

    rect(sl, 5.55, 1.55, 7.55, 0.38, NAVY)
    txt(sl, 'COMMON INTERVIEW QUESTIONS — MODEL ANSWER STRUCTURES', 5.6, 1.58, 7.4, 0.32,
        size=8.5, bold=True, color=WHITE)

    for i, (q, ans) in enumerate(questions):
        y = 1.98 + i * 0.79
        rect(sl, 5.55, y, 7.55, 0.75, GREEN if i % 2 == 0 else LGREENB)
        qcolor = WHITE if i % 2 == 0 else NAVY
        acolor = WHITE if i % 2 == 0 else NAVY
        txt(sl, f'Q{i+1}: {q}', 5.65, y+0.04, 7.3, 0.3, size=8.5, bold=True, color=qcolor)
        txt(sl, ans, 5.65, y+0.34, 7.3, 0.4, size=7.5, color=acolor, italic=True)

    return sl


def slide_checklist(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, LGREY)
    section_header(sl, 'Interview Prep', 'Key Facts to Nail + Red Flags to Avoid',
                   'Walk in knowing these — walk out having impressed')
    nav_bar(sl, 7)

    musts = [
        'Know your exact SSE division & its lead projects by name',
        'Name-drop: Dogger Bank, Viking, Peterhead CCS, Pathway to 2030',
        'Cite the £33bn five-year plan — 80% goes to the electricity grid',
        '"Safety is SSE\'s number one value" — say this naturally',
        'Use "energy transition" and "net zero" confidently throughout',
        'SSE targets 16 GW renewable capacity by 2032 (fivefold increase)',
        'SSE delivers 25%+ of UK\'s 2030 40GW offshore wind target',
        'Grid intermittency makes flexible generation + storage essential',
        'Ask: "What would you expect from me in the first 90 days?"',
        '83,000+ jobs supported — SSE is a major UK economic engine',
    ]

    avoid = [
        'Not knowing which SSE business unit you applied for',
        'Generic "I like renewable energy" without naming specific projects',
        'Confusing SSE plc with Stockholm School of Economics',
        'Bluffing on technical questions — say "I\'d approach it by..."',
        'Not preparing STAR answers before the 3-min think window',
        'Overlooking Safety as your #1 competency example',
    ]

    # Must-knows
    rect(sl, 0.3, 1.55, 6.5, 0.45, GREEN)
    txt(sl, '✓  KEY FACTS TO DROP IN THE INTERVIEW', 0.4, 1.58, 6.3, 0.38, size=10, bold=True, color=WHITE)
    rect(sl, 0.3, 2.03, 6.5, 5.0, WHITE)
    for i, item in enumerate(musts):
        y = 2.1 + i * 0.48
        rect(sl, 0.35, y+0.04, 0.35, 0.34, GREEN)
        txt(sl, '✓', 0.35, y+0.06, 0.35, 0.3, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, item, 0.8, y+0.05, 5.85, 0.38, size=8.8, color=NAVY)

    # Avoid
    rect(sl, 7.0, 1.55, 6.1, 0.45, NAVY)
    txt(sl, '✗  RED FLAGS TO AVOID', 7.1, 1.58, 5.9, 0.38, size=10, bold=True, color=WHITE)
    rect(sl, 7.0, 2.03, 6.1, 4.3, WHITE)
    for i, item in enumerate(avoid):
        y = 2.1 + i * 0.7
        rect(sl, 7.05, y+0.05, 0.35, 0.52, RGBColor(0xE5,0x3E,0x3E))
        txt(sl, '✗', 7.05, y+0.08, 0.35, 0.4, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, item, 7.5, y+0.06, 5.5, 0.55, size=9, color=NAVY)

    return sl


def slide_closing(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, NAVY)
    rect(sl, 0, 0, 0.18, 7.5, GREEN)
    rect(sl, 0, 3.45, 13.33, 0.08, GREEN)

    txt(sl, 'You\'re Ready.', 0.5, 0.6, 12, 1.2, size=52, bold=True, color=WHITE)
    txt(sl, 'SSE Interview Preparation — Complete', 0.5, 1.65, 12, 0.55,
        size=18, color=GREEN, bold=True)

    summary = [
        ('Company', 'FTSE 100 | Perth | £20.93bn | Est. 1998'),
        ('Investment', '£33bn 5-year plan | 80% grid | £7M/day'),
        ('Projects', 'Dogger Bank · Viking · Peterhead CCS · Pathway to 2030'),
        ('Values', 'Safety (first) · Service · Efficiency · Sustainability · Excellence · Teamwork'),
        ('Process', 'Application → SHL Tests → Video Interview → Assessment Centre'),
        ('Your Edge', 'EE Masters: power systems, HV design, protection, power electronics'),
    ]
    for i, (label, detail) in enumerate(summary):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.4
        y = 3.8 + row * 1.08
        rect(sl, x, y, 0.7, 0.75, GREEN)
        txt(sl, label, x+0.05, y+0.18, 0.62, 0.38, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, x+0.75, y, 5.5, 0.75, RGBColor(0x22,0x3A,0x70))
        txt(sl, detail, x+0.85, y+0.15, 5.3, 0.45, size=9.5, color=WHITE)

    txt(sl, 'Good luck — the energy sector needs engineers who understand both the grid and the transition.',
        0.5, 6.85, 12.5, 0.45, size=10, color=RGBColor(0xAA,0xBB,0xDD), italic=True)
    return sl


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == '__main__':
    prs = new_prs()

    print("Building slides...")
    slide_title(prs);               print("  ✓ Title slide")
    slide_company_overview(prs);    print("  ✓ Company Overview")
    slide_business_divisions(prs);  print("  ✓ Business Divisions")
    slide_key_projects(prs);        print("  ✓ Key Projects")
    slide_investment(prs);          print("  ✓ Investment Plan (donut chart)")
    slide_capacity_timeline(prs);   print("  ✓ Capacity & Timeline (bar + timeline chart)")
    slide_uk_impact(prs);           print("  ✓ UK Economic Impact (bar chart)")
    slide_values(prs);              print("  ✓ Company Values (radar chart)")
    slide_ee_masters(prs);          print("  ✓ EE Masters Mapping (horizontal bar chart)")
    slide_interview_process(prs);   print("  ✓ Recruitment Process (funnel)")
    slide_star_questions(prs);      print("  ✓ STAR + Interview Questions")
    slide_checklist(prs);           print("  ✓ Checklist + Red Flags")
    slide_closing(prs);             print("  ✓ Closing Summary")

    out = '/home/user/data-scientist-ai-era/SSE_Interview_Preparation.pptx'
    prs.save(out)
    print(f"\n✅ Saved: {out}")
